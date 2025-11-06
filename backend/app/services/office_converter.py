"""
Office文档转换服务
将Office文档转换为PDF以便在浏览器中预览
"""

import os
import io
import aiofiles
from typing import Optional, Dict, Any
from pathlib import Path
import subprocess
import tempfile

from docx import Document
from pptx import Presentation
import PyPDF2
from PIL import Image
import fitz  # PyMuPDF


class OfficeConverterService:
    """Office文档转换服务"""

    def __init__(self):
        self.temp_dir = tempfile.gettempdir()

    async def convert_to_pdf(self, file_path: str, output_path: str) -> Dict[str, Any]:
        """
        将Office文档转换为PDF

        Args:
            file_path: 源文件路径
            output_path: 输出PDF路径

        Returns:
            转换结果信息
        """
        try:
            file_ext = Path(file_path).suffix.lower()

            if file_ext == ".docx":
                return await self._convert_docx_to_pdf(file_path, output_path)
            elif file_ext in [".ppt", ".pptx"]:
                return await self._convert_ppt_to_pdf(file_path, output_path)
            else:
                raise ValueError(f"不支持的文档格式: {file_ext}")

        except Exception as e:
            return {"success": False, "error": str(e), "pdf_url": None}

    async def _convert_docx_to_pdf(self, docx_path: str, pdf_path: str) -> Dict[str, Any]:
        """将DOCX转换为PDF"""
        try:
            # 方法1: 使用LibreOffice命令行工具（推荐）
            if await self._has_libreoffice():
                return await self._convert_with_libreoffice(docx_path, pdf_path)

            # 方法2: 使用python-docx提取内容并生成简化PDF
            return await self._convert_docx_content_to_pdf(docx_path, pdf_path)

        except Exception as e:
            return {"success": False, "error": f"DOCX转换失败: {str(e)}", "pdf_url": None}

    async def _convert_ppt_to_pdf(self, ppt_path: str, pdf_path: str) -> Dict[str, Any]:
        """将PPT/PPTX转换为PDF"""
        try:
            # 方法1: 使用LibreOffice命令行工具（推荐）
            if await self._has_libreoffice():
                return await self._convert_with_libreoffice(ppt_path, pdf_path)

            # 方法2: 使用python-pptx提取内容并生成简化PDF
            return await self._convert_ppt_content_to_pdf(ppt_path, pdf_path)

        except Exception as e:
            return {"success": False, "error": f"PPT转换失败: {str(e)}", "pdf_url": None}

    async def _has_libreoffice(self) -> bool:
        """检查是否安装了LibreOffice"""
        try:
            result = subprocess.run(
                ["libreoffice", "--version"], capture_output=True, text=True, timeout=5
            )
            return result.returncode == 0
        except:
            return False

    async def _convert_with_libreoffice(self, input_path: str, output_path: str) -> Dict[str, Any]:
        """使用LibreOffice转换文档"""
        try:
            output_dir = os.path.dirname(output_path)

            # 使用LibreOffice命令行转换
            cmd = [
                "libreoffice",
                "--headless",
                "--convert-to",
                "pdf",
                "--outdir",
                output_dir,
                input_path,
            ]

            result = subprocess.run(
                cmd, capture_output=True, text=True, timeout=90
            )  # 增加到90秒，适应大文件转换

            if result.returncode == 0:
                # 获取生成的PDF文件名
                input_name = Path(input_path).stem
                generated_pdf = os.path.join(output_dir, f"{input_name}.pdf")

                if os.path.exists(generated_pdf) and generated_pdf != output_path:
                    # 重命名到目标路径
                    os.rename(generated_pdf, output_path)

                return {
                    "success": True,
                    "error": None,
                    "pdf_url": output_path,
                    "method": "libreoffice",
                }
            else:
                return {
                    "success": False,
                    "error": f"LibreOffice转换失败: {result.stderr}",
                    "pdf_url": None,
                }

        except Exception as e:
            return {"success": False, "error": f"LibreOffice转换异常: {str(e)}", "pdf_url": None}

    async def _convert_docx_content_to_pdf(self, docx_path: str, pdf_path: str) -> Dict[str, Any]:
        """从DOCX提取内容并生成简化PDF"""
        try:
            # 读取DOCX文档
            doc = Document(docx_path)

            # 创建PDF文档
            pdf_doc = fitz.open()

            # 提取段落内容
            y_position = 50
            page_margin = 50
            line_height = 20
            max_width = 500

            for paragraph in doc.paragraphs:
                if paragraph.text.strip():
                    # 检查是否需要新页面
                    if y_position > 750:  # 页面底部
                        pdf_doc.new_page()
                        y_position = 50

                    # 处理文本格式
                    text = paragraph.text.strip()

                    # 检查是否是标题（通过样式判断）
                    is_heading = False
                    if paragraph.style.name.startswith("Heading"):
                        is_heading = True
                        font_size = 16
                        font_color = (0, 0, 0)  # 黑色
                    else:
                        font_size = 12
                        font_color = (0, 0, 0)  # 黑色

                    # 添加文本到PDF
                    try:
                        pdf_doc[-1].insert_text(
                            (page_margin, y_position), text, fontsize=font_size, color=font_color
                        )
                        y_position += line_height + (10 if is_heading else 0)
                    except Exception as text_error:
                        print(f"Error inserting text: {text_error}")
                        # 如果插入失败，尝试简单的文本插入
                        try:
                            pdf_doc[-1].insert_text((page_margin, y_position), text, fontsize=12)
                            y_position += line_height
                        except:
                            continue

            # 如果没有内容，添加提示信息
            if len(pdf_doc) == 0 or not any(page.get_text().strip() for page in pdf_doc):
                page = pdf_doc.new_page()
                page.insert_text((50, 100), "文档内容提取中...", fontsize=14)
                page.insert_text((50, 130), "请使用其他预览方式查看完整内容", fontsize=12)

            # 保存PDF
            pdf_doc.save(pdf_path)
            pdf_doc.close()

            return {
                "success": True,
                "error": None,
                "pdf_url": pdf_path,
                "method": "content_extraction",
            }

        except Exception as e:
            return {"success": False, "error": f"内容提取转换失败: {str(e)}", "pdf_url": None}

    async def _convert_ppt_content_to_pdf(self, ppt_path: str, pdf_path: str) -> Dict[str, Any]:
        """从PPT/PPTX提取内容并生成简化PDF"""
        try:
            # 读取PPT文档
            prs = Presentation(ppt_path)

            # 创建PDF文档
            pdf_doc = fitz.open()

            for slide_num, slide in enumerate(prs.slides):
                page = pdf_doc.new_page()

                # 添加幻灯片标题
                page.insert_text((50, 50), f"幻灯片 {slide_num + 1}", fontsize=16, color=(0, 0, 0))

                # 提取文本框内容
                y_pos = 100
                content_found = False

                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        content_found = True
                        text = shape.text.strip()

                        # 检查文本长度，如果太长则分行
                        if len(text) > 80:
                            # 简单的文本换行
                            words = text.split()
                            lines = []
                            current_line = ""

                            for word in words:
                                if len(current_line + " " + word) <= 80:
                                    current_line += " " + word if current_line else word
                                else:
                                    if current_line:
                                        lines.append(current_line)
                                    current_line = word

                            if current_line:
                                lines.append(current_line)

                            for line in lines:
                                if y_pos > 750:  # 页面底部，创建新页面
                                    page = pdf_doc.new_page()
                                    page.insert_text(
                                        (50, 50),
                                        f"幻灯片 {slide_num + 1} (续)",
                                        fontsize=14,
                                        color=(0, 0, 0),
                                    )
                                    y_pos = 100

                                page.insert_text((50, y_pos), line, fontsize=12, color=(0, 0, 0))
                                y_pos += 25
                        else:
                            if y_pos > 750:  # 页面底部
                                page = pdf_doc.new_page()
                                page.insert_text(
                                    (50, 50),
                                    f"幻灯片 {slide_num + 1} (续)",
                                    fontsize=14,
                                    color=(0, 0, 0),
                                )
                                y_pos = 100

                            page.insert_text((50, y_pos), text, fontsize=12, color=(0, 0, 0))
                            y_pos += 25

                # 如果没有找到文本内容，添加提示
                if not content_found:
                    page.insert_text(
                        (50, 100), "此幻灯片无文本内容", fontsize=12, color=(100, 100, 100)
                    )
                    page.insert_text(
                        (50, 130), "可能包含图片或其他媒体内容", fontsize=10, color=(100, 100, 100)
                    )

            # 如果没有幻灯片，添加提示信息
            if len(pdf_doc) == 0:
                page = pdf_doc.new_page()
                page.insert_text((50, 100), "演示文稿内容提取中...", fontsize=14)
                page.insert_text((50, 130), "请使用其他预览方式查看完整内容", fontsize=12)

            # 保存PDF
            pdf_doc.save(pdf_path)
            pdf_doc.close()

            return {
                "success": True,
                "error": None,
                "pdf_url": pdf_path,
                "method": "content_extraction",
            }

        except Exception as e:
            return {"success": False, "error": f"PPT内容提取转换失败: {str(e)}", "pdf_url": None}

    async def get_converted_pdf_url(self, original_file_url: str) -> Optional[str]:
        """
        获取Office文档的转换PDF URL

        Args:
            original_file_url: 原始文件URL

        Returns:
            转换后的PDF URL，如果转换失败则返回None
        """
        try:
            # 从URL提取文件路径
            if not original_file_url.startswith("/uploads/resources/"):
                return None

            # 构建文件路径
            file_path = original_file_url.replace("/uploads/resources/", "storage/resources/")
            if not os.path.exists(file_path):
                return None

            # 检查是否已经转换过
            pdf_path = file_path.rsplit(".", 1)[0] + "_converted.pdf"
            if os.path.exists(pdf_path):
                return pdf_path.replace("storage/resources/", "/uploads/resources/")

            # 执行转换
            result = await self.convert_to_pdf(file_path, pdf_path)

            if result["success"]:
                print(f"Office文档转换成功，使用方法: {result.get('method', 'unknown')}")
                return pdf_path.replace("storage/resources/", "/uploads/resources/")
            else:
                print(f"Office文档转换失败: {result['error']}")
                return None

        except Exception as e:
            print(f"获取转换PDF URL失败: {str(e)}")
            return None


# 单例实例
office_converter_service = OfficeConverterService()
